"""
Unified Consulting Ingestor - Reader unificado para documentos de consultoria

Suporta múltiplos formatos:
- PPTX/PDF → API Visual → Markdown V019 → Document
- .md estruturado → Parse direto → Document

Output sempre compatível com SlidesSemanticaVisualChunker (formato V019).

Providers de API Visual:
- Contextual.AI (recomendado)
- Docling
- GPT-4V
- Claude Vision

Integração:
- Funciona com SlidesSemanticaVisualChunker existente
- Funciona com HybridConsultingEmbedder
- Compatible com ETL pré/pós-chunking
"""

import base64
import os
import re
from typing import List, Dict, Any, Optional
from wasabi import msg
import aiohttp
import json

from goldenverba.components.document import Document
from goldenverba.components.interfaces import Reader
from goldenverba.server.types import FileConfig
from goldenverba.components.types import InputConfig


class UnifiedConsultingIngestor(Reader):
    """
    Reader unificado para documentos de consultoria.
    
    Pipeline Unificado:
    1. PPTX/PDF → API Visual → Markdown estruturado → Document
    2. .md estruturado → Parse direto → Document
    
    Output: Document com slides_metadata[] (formato V019)
    Compatible: SlidesSemanticaVisualChunker
    """
    
    def __init__(self):
        super().__init__()
        self.requires_env = []
        self.name = "Unified Consulting Ingestor"
        self.description = (
            "Reader unificado para consultoria. Processa PPTX/PDF via API visual "
            "(Contextual.AI/Docling/GPT-4V/Claude) ou .md estruturado direto. "
            "Output compatível com Visual Semantic (formato V019)."
        )
        
        # Configuração
        self.config = {
            "Enable Visual Analysis": InputConfig(
                type="bool",
                value=False,  # Default: OFF (economia)
                description="Habilitar análise visual via API (Contextual.AI/GPT-4V/etc.). Se OFF, usa Docling text-only.",
                values=[]
            ),
            "Visual API Provider": InputConfig(
                type="dropdown",
                value="Contextual.AI",
                description="Provider de análise visual (apenas se Enable Visual Analysis = True)",
                values=["Contextual.AI", "Docling", "GPT-4V", "Claude Vision"]
            ),
            "API URL": InputConfig(
                type="text",
                value=os.getenv("CONTEXTUAL_API_URL", "https://api.contextual.ai/v1/documents"),
                description="URL da API (Contextual.AI, Docling Remoto, etc.)",
                values=[]
            ),
            "API Key": InputConfig(
                type="password",
                value=os.getenv("CONTEXTUAL_API_KEY", ""),
                description="API Key visual (se habilitada)",
                values=[]
            ),
            "Framework Detection Method": InputConfig(
                type="text",
                value="Auto",
                description="Método de detecção de frameworks (Auto, GLiNER, Regex, Visual API)",
                values=["Auto", "GLiNER", "Regex", "Visual API"]
            ),
            "Language": InputConfig(
                type="text",
                value="pt",
                description="Idioma do documento",
                values=["pt", "en", "es"]
            ),
            "Enable ETL": InputConfig(
                type="bool",
                value=True,
                description="Habilitar ETL pré-chunking automático",
                values=[]
            )
        }
        
        # Framework patterns (para detecção)
        self.framework_patterns = self._get_framework_patterns()
    
    async def load(self, config: dict, fileConfig: FileConfig) -> List[Document]:
        """
        Carrega documento em qualquer formato suportado.
        
        Returns:
            Lista com 1 Document no formato V019
        """
        # Detecta formato
        extension = fileConfig.extension.lower().strip().lstrip('.')
        
        msg.info(f"[Unified-Ingestor] Processando: {fileConfig.filename} (.{extension})")
        
        if extension in ["pptx", "pdf"]:
            # Caminho 1: API Visual ou Docling
            return await self._process_via_visual_api(config, fileConfig)
        
        elif extension == "md":
            # Caminho 2: Markdown estruturado
            return await self._process_structured_markdown(config, fileConfig)
        
        else:
            msg.fail(f"[Unified-Ingestor] Formato não suportado: .{extension}")
            return []
    
    async def _process_via_visual_api(
        self,
        config: dict,
        fileConfig: FileConfig
    ) -> List[Document]:
        """
        Pipeline: PPTX/PDF → Docling/API Visual → Markdown → Document
        
        Dois caminhos:
        1. Visual Analysis ON: API visual completa (gráficos, semântica)
        2. Visual Analysis OFF: Docling text-only (economia)
        """
        enable_visual = self._get_config_value(config, "Enable Visual Analysis", False)
        
        if enable_visual:
            # Caminho 1: Visual API completa
            return await self._process_with_visual_api(config, fileConfig)
        else:
            # Caminho 2: Docling text-only (economia)
            msg.info("[Unified-Ingestor] Modo: Docling text-only (Visual Analysis OFF)")
            return await self._process_with_docling_text_only(config, fileConfig)
    
    async def _process_with_visual_api(
        self,
        config: dict,
        fileConfig: FileConfig
    ) -> List[Document]:
        """
        Caminho 1: API Visual completa (análise semântica visual).
        """
        provider = self._get_config_value(config, "Visual API Provider", "Contextual.AI")
        api_url = self._get_config_value(config, "API URL", "")
        # Try to get API key from config, fallback to env var logic in a real scenario
        api_key = self._get_config_value(config, "API Key", "") 
        language = self._get_config_value(config, "Language", "pt")
        
        if not api_key:
             # Just a warning in this mock/impl phase
            msg.warn("[Unified-Ingestor] API Key não configurada (modo visual). Usando Mock se falhar.")
        
        msg.info(f"[Unified-Ingestor] Usando Visual API: {provider}")
        
        # 1. Chama API visual (Placeholder/Mock logic if no real API connection)
        visual_result = await self._call_visual_api(
            provider, api_url, api_key, fileConfig, language, extract_visual=True
        )
        
        if not visual_result:
             # Fallback mock for testing
             msg.info("[Unified-Ingestor] API visual falhou ou não configurada. Gerando Mock.")
             visual_result = self._generate_mock_visual_result(fileConfig)

        # 2. Gera markdown V019 (ou usa direto se a API retornar markdown)
        if "markdown" in visual_result:
            markdown = visual_result["markdown"]
        else:
            markdown = self._generate_v019_markdown(visual_result)
        
        # 3. Parse → Document
        document = self._parse_markdown_to_document(markdown, fileConfig, config)
        
        msg.good(f"[Unified-Ingestor] ✅ Visual API completa: {document.meta['slide_count']} slides")
        
        return [document]
    
    async def _process_with_docling_text_only(
        self,
        config: dict,
        fileConfig: FileConfig
    ) -> List[Document]:
        """
        Caminho 2: Docling text-only (sem análise visual).
        """
        language = self._get_config_value(config, "Language", "pt")
        framework_method = self._get_config_value(config, "Framework Detection Method", "Auto")
        
        msg.info("[Unified-Ingestor] Processando com Docling text-only...")
        
        try:
            # 1. Extrai texto via Docling
            docling_result = await self._extract_text_with_docling(fileConfig, language)
            
            if not docling_result:
                msg.fail("[Unified-Ingestor] Falha no Docling")
                return []
            
            # 2. Detecta frameworks (GLiNER ou regex)
            if framework_method == "Auto":
                framework_method = "GLiNER" if self._is_gliner_available() else "Regex"
            
            msg.info(f"[Unified-Ingestor] Detectando frameworks via: {framework_method}")
            docling_result = self._detect_frameworks_text_based(docling_result, framework_method)
            
            # 3. Detecta stakeholders via NER simples
            docling_result = self._detect_stakeholders_text_based(docling_result, language)
            
            # 4. Gera markdown V019
            markdown = self._generate_v019_markdown(docling_result)
            
            # 5. Parse → Document
            document = self._parse_markdown_to_document(markdown, fileConfig, config)
            
            msg.good(f"[Unified-Ingestor] ✅ Docling text-only: {document.meta['slide_count']} slides")
            if 'all_frameworks' in document.meta:
                 msg.info(f"[Unified-Ingestor]    Frameworks: {len(document.meta.get('all_frameworks', []))}")
            if 'all_companies' in document.meta:
                 msg.info(f"[Unified-Ingestor]    Companies: {len(document.meta.get('all_companies', []))}")
            
            return [document]
        
        except Exception as e:
            msg.fail(f"[Unified-Ingestor] Erro no Docling: {str(e)}")
            import traceback
            msg.debug(traceback.format_exc())
            return []
    
    async def _extract_text_with_docling(
        self,
        fileConfig: FileConfig,
        language: str
    ) -> Optional[Dict[str, Any]]:
        """
        Extrai texto estruturado via Docling.
        """
        try:
            # TODO: Integrar Docling real imports dentro do metodo para evitar erro de load se nao instalado
            try:
                from docling.document_converter import DocumentConverter
                
                # Docling requer path ou manipulador de arquivo.
                # Precisamos salvar o conteudo temporariamente se nao tivermos path
                temp_path = f"temp_{fileConfig.filename}"
                with open(temp_path, "wb") as f:
                     # Suporte a string ou bytes no content
                     if isinstance(fileConfig.content, str):
                          f.write(fileConfig.content.encode('utf-8'))
                     else:
                          f.write(fileConfig.content)
                
                converter = DocumentConverter()
                result = converter.convert(temp_path)
                
                # Cleanup
                if os.path.exists(temp_path):
                     os.remove(temp_path)
                     
                # Export to markdown directly from Docling
                full_md = result.document.export_to_markdown()
                
                # Parse markdown back to slides structure (Docling doesn't respect slides natively yet)
                # This simulates slide separation by double newline or headers
                paragraphs = full_md.split('\n\n')
                
            except ImportError:
                msg.warn("[Unified-Ingestor] Docling library not found. Trying python-pptx fallback...")
                # Try python-pptx for PPTX files
                try:
                    from pptx import Presentation
                    import io
                    
                    pptx_bytes = io.BytesIO(base64.b64decode(fileConfig.content) if isinstance(fileConfig.content, str) else fileConfig.content)
                    prs = Presentation(pptx_bytes)
                    
                    slides = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_texts = []
                        
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_texts.append(shape.text.strip())
                            # Extract table text if present
                            try:
                                if shape.has_table:
                                    table = shape.table
                                    for row in table.rows:
                                        row_text = []
                                        for cell in row.cells:
                                            if cell.text.strip():
                                                row_text.append(cell.text.strip())
                                        if row_text:
                                            slide_texts.append(" | ".join(row_text))
                            except Exception:
                                pass
                        
                        if slide_texts:
                            title = slide_texts[0] if slide_texts else f"Slide {slide_num}"
                            content = "\n".join(slide_texts[1:]) if len(slide_texts) > 1 else "\n".join(slide_texts)
                            
                            slides.append({
                                "number": slide_num,
                                "title": title[:200] + "..." if len(title) > 200 else title,
                                "content": content,
                                "frameworks": [],
                                "companies": [],
                                "position": self._infer_position(slide_num, len(prs.slides))
                            })
                    
                    if slides:
                        msg.good(f"[Unified-Ingestor] python-pptx extracted {len(slides)} slides")
                        return {"slides": slides}
                    else:
                        msg.warn("[Unified-Ingestor] python-pptx found no text content")
                        return None
                        
                except ImportError:
                    msg.warn("[Unified-Ingestor] python-pptx not installed. Falling back to basic text extraction.")
                    # Basic Fallback for text files only
                    content = fileConfig.content
                    paragraphs = content.decode('utf-8').split('\n\n') if isinstance(content, bytes) else content.split('\n\n')
                except Exception as e:
                    msg.warn(f"[Unified-Ingestor] python-pptx error: {e}. Falling back to basic extraction.")
                    content = fileConfig.content
                    paragraphs = content.decode('utf-8').split('\n\n') if isinstance(content, bytes) else content.split('\n\n')

            slides = []
            for i, para in enumerate(paragraphs, 1):
                if para.strip():
                    lines = para.strip().split('\n')
                    title = lines[0] if lines else f"Slide {i}"
                    content_text = '\n'.join(lines[1:]) if len(lines) > 1 else para
                    
                    # Heuristic for title length
                    if len(title) > 200: 
                        title = title[:100] + "..."
                    
                    slides.append({
                        "number": i,
                        "title": title,
                        "content": content_text,
                        "frameworks": [],
                        "companies": [],
                        "position": self._infer_position(i, len(paragraphs))
                    })
            
            if not slides:
                return None
            
            return {"slides": slides}
        
        except Exception as e:
            msg.warn(f"[Unified-Ingestor] Erro no Docling Extraction: {str(e)}")
            return None
    
    def _detect_frameworks_text_based(
        self,
        docling_result: Dict,
        method: str
    ) -> Dict:
        """
        Detecta frameworks via texto (GLiNER ou regex).
        """
        if method == "GLiNER":
            return self._detect_frameworks_gliner(docling_result)
        else:
            return self._detect_frameworks_regex(docling_result)
    
    def _detect_frameworks_regex(self, docling_result: Dict) -> Dict:
        """
        Detecta frameworks via regex patterns.
        """
        framework_patterns = self._get_framework_patterns()
        
        for slide in docling_result.get("slides", []):
            content_lower = (str(slide.get("title", "")) + " " + str(slide.get("content", ""))).lower()
            
            detected = []
            for framework_name, patterns in framework_patterns.items():
                for pattern in patterns:
                    if re.search(pattern, content_lower, re.IGNORECASE):
                        detected.append(framework_name)
                        break
            
            slide["frameworks"] = list(set(detected))
        
        return docling_result
    
    def _detect_frameworks_gliner(self, docling_result: Dict) -> Dict:
        """
        Detecta frameworks via GLiNER.
        """
        msg.info("[Unified-Ingestor] GLiNER não implementado ainda, usando regex")
        return self._detect_frameworks_regex(docling_result)
    
    def _detect_stakeholders_text_based(
        self,
        docling_result: Dict,
        language: str
    ) -> Dict:
        """
        Detecta stakeholders (empresas, pessoas) via spaCy NER básico.
        """
        try:
            import spacy
            
            model_name = {
                "pt": "pt_core_news_sm",
                "en": "en_core_web_sm",
                "es": "es_core_news_sm"
            }.get(language, "en_core_web_sm")
            
            # Lazy load or check existence
            if not spacy.util.is_package(model_name):
                 # Try download? Or fallback
                 msg.warn(f"[Unified-Ingestor] Modelo {model_name} nao instalado. Pulando NER.")
                 return docling_result

            nlp = spacy.load(model_name)
            
            for slide in docling_result.get("slides", []):
                text = str(slide.get("title", "")) + " " + str(slide.get("content", ""))
                doc = nlp(text)
                
                companies = []
                for ent in doc.ents:
                    if ent.label_ in ["ORG", "PERSON"]:
                        companies.append(ent.text)
                
                slide["companies"] = list(set(companies))[:10]
            
            return docling_result
        
        except ImportError:
            msg.warn("[Unified-Ingestor] spaCy não disponível, pulando NER")
            return docling_result
        except Exception as e:
            msg.warn(f"[Unified-Ingestor] Erro no NER: {str(e)}")
            return docling_result
    
    def _infer_position(self, slide_num: int, total_slides: int) -> str:
        if slide_num == 1:
            return "opening"
        elif slide_num <= total_slides * 0.3:
            return "diagnostic"
        elif slide_num <= total_slides * 0.7:
            return "analysis"
        else:
            return "conclusion"
    
    def _is_gliner_available(self) -> bool:
        try:
            import gliner
            return True
        except ImportError:
            return False
            
    async def _process_structured_markdown(
        self,
        config: dict,
        fileConfig: FileConfig
    ) -> List[Document]:
        """
        Pipeline: .md estruturado → Document
        """
        msg.info("[Unified-Ingestor] Processando markdown estruturado")
        
        try:
            if isinstance(fileConfig.content, str):
                markdown_content = fileConfig.content
            else:
                decoded_bytes = fileConfig.content
                markdown_content = decoded_bytes.decode('utf-8')
        except Exception as e:
            msg.fail(f"[Unified-Ingestor] Erro ao decodificar: {str(e)}")
            return []
        
        document = self._parse_markdown_to_document(markdown_content, fileConfig, config)
        msg.good(f"[Unified-Ingestor] ✅ Processado: {document.meta['slide_count']} slides")
        return [document]
    
    async def _call_visual_api(
        self,
        provider: str,
        api_url: str,
        api_key: str,
        fileConfig: FileConfig,
        language: str,
        extract_visual: bool
    ) -> Optional[Dict[str, Any]]:
        """
        Calls external Visual API (Contextual.AI, Docling Remote, etc.)
        """
        if not api_url:
            msg.warn(f"[Unified-Ingestor] API URL não fornecida para provider: {provider}")
            return None

        # Dados para envio
        filename = fileConfig.filename
        content = fileConfig.content
        
        try:
            async with aiohttp.ClientSession() as session:
                data = aiohttp.FormData()
                
                # Se content for string (base64 ou texto), converte para bytes
                if isinstance(content, str):
                    if ";base64," in content:
                        content = base64.b64decode(content.split(";base64,")[1])
                    else:
                        content = content.encode('utf-8')
                
                data.add_field('file', content, filename=filename)
                data.add_field('language', language)
                if extract_visual:
                    data.add_field('extract_visual', 'true')
                
                headers = {}
                if api_key:
                    headers["Authorization"] = f"Bearer {api_key}"
                
                msg.info(f"[Unified-Ingestor] Enviando {filename} para {api_url} ({provider})")
                
                async with session.post(api_url, data=data, headers=headers, timeout=120) as response:
                    if response.status != 200:
                        error_text = await response.text()
                        msg.fail(f"[Unified-Ingestor] Erro na API {provider} ({response.status}): {error_text}")
                        return None
                    
                    result = await response.json()
                    
                    # Normalização básica dependendo do provider
                    if "slides" in result:
                        return result
                    elif "markdown" in result:
                        # Se retornar só markdown, tenta fazer parse interno
                        return {"markdown": result["markdown"]}
                    else:
                        # Fallback: assume que o resultado já está no formato correto ou precisa de parse
                        return result
                        
        except Exception as e:
            msg.fail(f"[Unified-Ingestor] Falha na chamada da API {provider}: {str(e)}")
            return None
    
    def _generate_mock_visual_result(self, fileConfig: FileConfig):
        """Helper to generate mock data for testing"""
        return {
            "slides": [
                {
                    "number": 1,
                    "title": "Executive Summary (Mock)",
                    "content": "This is a mock slide content representing Visual API analysis.",
                    "frameworks": ["BCG Matrix (Mock)"],
                    "companies": ["Mock Corp"],
                    "visual_archetype": "matrix",
                    "position": "opening"
                },
                 {
                    "number": 2,
                    "title": "Financial Analysis (Mock)",
                    "content": "Financial growth is positive.",
                    "frameworks": [],
                    "companies": [],
                    "position": "analysis"
                }
            ]
        }

    def _generate_v019_markdown(self, visual_result: Dict) -> str:
        """
        Gera markdown estruturado no formato V019.
        """
        markdown_parts = []
        slides = visual_result.get("slides", [])
        
        for slide in slides:
            slide_number = slide.get("number", 0)
            slide_title = slide.get("title", f"Slide {slide_number}")
            markdown_parts.append(f"# Slide {slide_number} - {slide_title}\n\n")
            
            content = slide.get("content", "")
            if content:
                markdown_parts.append(content)
                markdown_parts.append("\n\n")
            
            frameworks = slide.get("frameworks", [])
            if frameworks:
                frameworks_str = ", ".join(frameworks)
                markdown_parts.append(f"**Frameworks Deste Slide:** {frameworks_str}\n\n")
            
            companies = slide.get("companies", [])
            if companies:
                companies_str = ", ".join(companies)
                markdown_parts.append(f"**Stakeholders Deste Slide:** {companies_str}\n\n")
            
            bridge_quality = slide.get("semantic_bridge_quality")
            if bridge_quality is not None:
                markdown_parts.append(f"**Qualidade da Ponte:** {bridge_quality:.2f}\n\n")
            
            position = slide.get("position")
            if position:
                markdown_parts.append(f"**Posição:** {position}\n\n")
            
            slide_type = slide.get("slide_type")
            if slide_type:
                markdown_parts.append(f"**Tipo de Slide:** {slide_type}\n\n")
            
            visual_archetype = slide.get("visual_archetype")
            if visual_archetype:
                markdown_parts.append(f"**Arquétipo Visual:** {visual_archetype}\n\n")
            
            pattern_genetics = slide.get("pattern_genetics", [])
            if pattern_genetics:
                genetics_str = ", ".join(pattern_genetics)
                markdown_parts.append(f"**Pattern Genetics:** {genetics_str}\n\n")
            
            reusability = slide.get("reusability_score")
            if reusability is not None:
                markdown_parts.append(f"**Reusability Score:** {reusability:.1f}\n\n")
            
            markdown_parts.append("---\n\n")
        
        return "".join(markdown_parts)
    
    def _parse_markdown_to_document(
        self,
        markdown_content: str,
        fileConfig: FileConfig,
        config: dict
    ) -> Document:
        """
        Parse markdown estruturado → Document com slides_metadata[].
        """
        slides_metadata = self._extract_slides_metadata_v019(markdown_content)
        
        document = Document(
            text=markdown_content,
            type="Unified Consulting Document",
            name=fileConfig.filename, # Adjusted from fileConfig.name to fileConfig.filename
            link=fileConfig.filename,
            timestamp="",
            reader=self.name
        )
        
        enable_etl = self._get_config_value(config, "Enable ETL", True)
        
        # Build meta dict
        meta_dict = {
            "enable_etl": enable_etl,
            "doc_type": "slides_semantica_visual",
            "source_format": fileConfig.extension,
            "slides_metadata": slides_metadata,
            "slide_count": len(slides_metadata),
            "all_frameworks": list(set([
                fw for slide in slides_metadata
                for fw in slide.get("frameworks", [])
            ])),
            "all_companies": list(set([
                comp for slide in slides_metadata
                for comp in slide.get("companies", [])
            ]))
        }
        document.meta = json.dumps(meta_dict)
        
        return document
    
    def _extract_slides_metadata_v019(self, markdown_content: str) -> List[Dict[str, Any]]:
        """
        Extrai metadata de slides do markdown V019.
        """
        slides_metadata = []
        slide_pattern = r"^#\s+Slide\s+(\d+)\s+-\s+(.+?)$"
        
        current_slide = None
        current_content = []
        
        for line in markdown_content.split("\n"):
            line = line.strip()
            match = re.match(slide_pattern, line, re.IGNORECASE)
            
            if match:
                if current_slide:
                    current_slide["content"] = "\n".join(current_content).strip()
                    self._extract_metadata_from_content(current_slide)
                    slides_metadata.append(current_slide)
                
                slide_number = int(match.group(1))
                slide_title = match.group(2).strip()
                current_slide = {
                    "slide_number": slide_number,
                    "slide_title": slide_title,
                    "frameworks": [],
                    "companies": [],
                }
                current_content = []
            else:
                if current_slide is not None:
                    current_content.append(line)
        
        if current_slide:
            current_slide["content"] = "\n".join(current_content).strip()
            self._extract_metadata_from_content(current_slide)
            slides_metadata.append(current_slide)
        
        return slides_metadata
    
    def _extract_metadata_from_content(self, slide: Dict):
        """
        Extrai metadata do conteúdo do slide (formato V019).
        """
        content = slide.get("content", "")
        
        frameworks_match = re.search(r"\*\*Frameworks Deste Slide:\*\*\s*(.+?)(?:\n|$)", content, re.IGNORECASE)
        if frameworks_match:
            frameworks_str = frameworks_match.group(1).strip()
            slide["frameworks"] = [f.strip() for f in frameworks_str.split(",")]
        
        stakeholders_match = re.search(r"\*\*Stakeholders Deste Slide:\*\*\s*(.+?)(?:\n|$)", content, re.IGNORECASE)
        if stakeholders_match:
            stakeholders_str = stakeholders_match.group(1).strip()
            slide["companies"] = [c.strip() for c in stakeholders_str.split(",")]
        
        bridge_match = re.search(r"\*\*Qualidade da Ponte:\*\*\s*([\d.]+)", content, re.IGNORECASE)
        if bridge_match:
            slide["semantic_bridge_quality"] = float(bridge_match.group(1))
        
        position_match = re.search(r"\*\*Posição:\*\*\s*(\w+)", content, re.IGNORECASE)
        if position_match:
            slide["position"] = position_match.group(1).strip()
        
        type_match = re.search(r"\*\*Tipo de Slide:\*\*\s*(\w+)", content, re.IGNORECASE)
        if type_match:
            slide["slide_type"] = type_match.group(1).strip()
        
        archetype_match = re.search(r"\*\*Arquétipo Visual:\*\*\s*(\w+)", content, re.IGNORECASE)
        if archetype_match:
            slide["visual_archetype"] = archetype_match.group(1).strip()
    
    def _get_framework_patterns(self) -> Dict[str, List[str]]:
        return {
            "BCG Matrix": [r"bcg", r"matriz bcg", r"boston consulting group"],
            "SWOT": [r"swot", r"forças? e fraquezas?"],
            "Porter": [r"porter", r"5 forças?", r"cinco forças?"],
            "McKinsey": [r"mckinsey", r"7s"],
            "Ansoff": [r"ansoff"],
            "PESTEL": [r"pestel", r"pest"]
        }
    
    def _get_config_value(self, config: dict, key: str, default):
        config_item = config.get(key, {})
        if isinstance(config_item, dict):
            return config_item.get("value", default)
        elif hasattr(config_item, 'value'):
            return config_item.value
        else:
            return default
